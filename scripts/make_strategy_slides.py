"""Generate vol2atlas_strategy.pptx in the style of methods_survey.pptx.

Run:  python scripts/make_strategy_slides.py
Out:  vol2atlas_strategy.pptx

Visual language (matches methods_survey.pptx):
  Background  #101014       Title text   #EEEEEE   Inter 44pt bold
  Card fill   #1C1C22       Subtitle     #9A9AA6   Inter 14pt
  Accent      #8AB4FF       Card title   #8AB4FF   Inter 17pt bold
  Warn        #F28B82       Body text    #EEEEEE   Inter 13pt
  Good        #81C995       Footer       #9A9AA6   Inter 10pt
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---- palette (extracted from methods_survey.pptx) ---------------------------
BG       = RGBColor(0x10, 0x10, 0x14)
CARD     = RGBColor(0x1C, 0x1C, 0x22)
CARD_HI  = RGBColor(0x24, 0x24, 0x2C)
INK      = RGBColor(0xEE, 0xEE, 0xEE)
INK_DIM  = RGBColor(0x9A, 0x9A, 0xA6)
ACCENT   = RGBColor(0x8A, 0xB4, 0xFF)   # blue
GOOD     = RGBColor(0x81, 0xC9, 0x95)   # green
WARN     = RGBColor(0xF2, 0x8B, 0x82)   # red
ATLAS    = RGBColor(0xFD, 0xD6, 0x63)   # gold
PURPLE   = RGBColor(0xC5, 0x8A, 0xF9)
HAIRLINE = RGBColor(0x33, 0x33, 0x3D)
FONT     = "Inter"

SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)


# ---- low-level helpers ------------------------------------------------------
def new_pres():
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6])   # Blank
    # full-bleed background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, left, top, width, height, body, *,
         size=13, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, italic=False,
         line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def rich_text(slide, left, top, width, height, paragraphs,
              size=13, color=INK, line_spacing=1.25,
              anchor=MSO_ANCHOR.TOP):
    """`paragraphs` is a list of paragraph specs. Each paragraph spec is
    a list of (text, kwargs) tuples — kwargs may set bold/color/size/font.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for pi, runs in enumerate(paragraphs):
        para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = line_spacing
        for txt, kw in runs:
            r = para.add_run()
            r.text = txt
            r.font.name  = kw.get("font",  FONT)
            r.font.size  = Pt(kw.get("size",  size))
            r.font.bold  = kw.get("bold",  False)
            r.font.italic = kw.get("italic", False)
            r.font.color.rgb = kw.get("color", color)
        para.space_after = Pt(kw.get("after", 4))
    return tb


def card(slide, left, top, width, height, *,
         fill=CARD, border=HAIRLINE, border_w=0.75, rounded=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    # small corner radius for cleaner look
    if rounded:
        sh.adjustments[0] = 0.06
    return sh


def chip(slide, left, top, width, height, label, *,
         fill=CARD_HI, text_color=ACCENT, size=10, bold=True):
    sh = card(slide, left, top, width, height,
              fill=fill, border=fill, rounded=True)
    sh.adjustments[0] = 0.5
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p_ = tf.paragraphs[0]; p_.alignment = PP_ALIGN.CENTER
    r = p_.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = text_color
    return sh


def arrow(slide, x1, y1, x2, y2, *, color=INK_DIM, weight=1.25, head=True):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    if head:
        from lxml import etree
        from pptx.oxml.ns import qn
        ln = line.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('h', 'med')
    return line


# ---- page furniture ---------------------------------------------------------
def title(slide, body):
    text(slide, Inches(0.55), Inches(0.4), Inches(12.3), Inches(0.85),
         body, size=32, bold=True, color=INK)


def subtitle(slide, body):
    text(slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(0.5),
         body, size=14, color=INK_DIM, italic=False)
    # thin accent underline
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(1.62), Inches(0.65), Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False


def footer(slide, n, total, *, kicker=""):
    text(slide, Inches(0.55), Inches(7.13), Inches(8), Inches(0.3),
         kicker, size=9, color=INK_DIM)
    text(slide, Inches(11.5), Inches(7.13), Inches(1.4), Inches(0.3),
         f"{n} / {total}", size=10, color=INK_DIM, align=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDES
# ============================================================================
def slide_title(p, n, total):
    s = add_slide(p)
    # large kicker
    text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         "STRATEGY  •  vol2atlas", size=14, bold=True, color=ACCENT)
    text(s, Inches(0.7), Inches(2.45), Inches(12), Inches(1.5),
         "Scalable multi-modal alignment",
         size=46, bold=True, color=INK)
    text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
         "Aligning µCT and EM volumes into Allen CCF — at GB and TB scales — through one pipeline.",
         size=18, color=INK_DIM)

    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(4.45), Inches(1.2), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False

    text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "A. Mittone  ·  APS / Argonne National Laboratory  ·  2026-05-22",
         size=12, color=INK_DIM)
    footer(s, n, total)


def slide_problem(p, n, total):
    s = add_slide(p)
    title(s, "The problem")
    subtitle(s, "vol2atlas already handles µCT → CCF. EM is the new modality — and it is multi-TB.")

    col_w = Inches(6.0); col_h = Inches(4.8); top = Inches(2.1)

    # WORKS
    c1 = card(s, Inches(0.55), top, col_w, col_h)
    text(s, Inches(0.85), top + Inches(0.25), col_w - Inches(0.5), Inches(0.45),
         "TODAY  ·  works", size=12, bold=True, color=GOOD)
    text(s, Inches(0.85), top + Inches(0.65), col_w - Inches(0.5), Inches(0.5),
         "µCT → Allen CCF", size=22, bold=True, color=INK)
    rich_text(s, Inches(0.85), top + Inches(1.25),
              col_w - Inches(0.6), col_h - Inches(1.5), [
        [("Sample.  ", {"bold": True, "color": ACCENT}),
         ("µCT volume, µm-class voxels, GB scale.", {})],
        [("Reference.  ", {"bold": True, "color": ACCENT}),
         ("Allen CCF via BrainGlobe (atlas name).", {})],
        [("Display.  ", {"bold": True, "color": ACCENT}),
         ("napari with strided preview ≤ 192³.", {})],
        [("Warp.  ", {"bold": True, "color": ACCENT}),
         ("alignFull = chunkwise rigid, TB-safe.", {})],
        [("Refinement.  ", {"bold": True, "color": ACCENT}),
         ("ants / elastix on the prealigned NIfTIs.", {})],
    ], size=13, line_spacing=1.4)

    # NEEDED
    c2 = card(s, Inches(6.8), top, col_w, col_h)
    text(s, Inches(7.1), top + Inches(0.25), col_w - Inches(0.5), Inches(0.45),
         "NEEDED  ·  new", size=12, bold=True, color=WARN)
    text(s, Inches(7.1), top + Inches(0.65), col_w - Inches(0.5), Inches(0.5),
         "EM → µCT (in CCF)", size=22, bold=True, color=INK)
    rich_text(s, Inches(7.1), top + Inches(1.25),
              col_w - Inches(0.6), col_h - Inches(1.5), [
        [("Sample.  ", {"bold": True, "color": ACCENT}),
         ("EM volume, nm-class voxels — 1000× finer.", {})],
        [("Reference.  ", {"bold": True, "color": ACCENT}),
         ("the already-aligned µCT, not an atlas.", {})],
        [("Scale.  ", {"bold": True, "color": ACCENT}),
         ("TB — RAM-bound warps are impossible.", {})],
        [("Pipeline.  ", {"bold": True, "color": ACCENT}),
         ("must accept any reference, not just BrainGlobe.", {})],
        [("Production warp.  ", {"bold": True, "color": ACCENT}),
         ("must compose rigid + TPS + (later) B-spline, chunkwise.", {})],
    ], size=13, line_spacing=1.4)

    footer(s, n, total, kicker="Problem statement")


def slide_chain(p, n, total):
    s = add_slide(p)
    title(s, "One pipeline, two stages — chained references")
    subtitle(s, "Stage B uses Stage A's output as its reference. Two independent runs. No two-stage state.")

    def node(left, top, w, h, label, *, color=ACCENT, body_size=14, dim=False):
        c = card(s, left, top, w, h,
                 fill=CARD if not dim else BG,
                 border=color, border_w=1.25)
        text(s, left, top + Inches(0.30), w, h - Inches(0.6),
             label, size=body_size, bold=True, color=INK if not dim else INK_DIM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return c

    # Stage A — top row
    yA = Inches(2.25); h = Inches(0.95)
    text(s, Inches(0.55), Inches(1.95), Inches(8), Inches(0.4),
         "STAGE A   ·   µCT  →  Allen CCF",
         size=12, bold=True, color=ACCENT)

    node(Inches(0.55), yA, Inches(2.2), h, "sample_uct.zarr", color=ACCENT)
    node(Inches(3.05), yA, Inches(4.0), h,
         "init → prealign → refine →\nlandmarks → alignFull",
         color=HAIRLINE, body_size=11)
    node(Inches(7.35), yA, Inches(2.4), h, "uct_in_ccf.zarr", color=ATLAS)
    node(Inches(10.05), yA, Inches(2.7), h,
         "reference:\nallen_mouse_25um",
         color=ATLAS, body_size=12, dim=True)

    arrow(s, Inches(2.75), yA + Emu(h.emu//2), Inches(3.05), yA + Emu(h.emu//2))
    arrow(s, Inches(7.05), yA + Emu(h.emu//2), Inches(7.35), yA + Emu(h.emu//2))
    arrow(s, Inches(10.05), yA + Emu(h.emu//2), Inches(9.75), yA + Emu(h.emu//2),
          color=ATLAS, weight=1.5)

    # Cross-stage handoff
    cx = Inches(8.55)
    arrow(s, cx, yA + h, cx, Inches(4.45), color=ATLAS, weight=2.0)
    text(s, Inches(8.85), Inches(3.55), Inches(4), Inches(0.4),
         "uct_in_ccf  →  reference for Stage B",
         size=11, bold=True, color=ATLAS, italic=True)

    # Stage B — bottom row
    yB = Inches(4.65)
    text(s, Inches(0.55), Inches(4.35), Inches(8), Inches(0.4),
         "STAGE B   ·   EM  →  µCT  (already in CCF)",
         size=12, bold=True, color=PURPLE)

    node(Inches(0.55), yB, Inches(2.2), h, "sample_em.zarr", color=PURPLE)
    node(Inches(3.05), yB, Inches(4.0), h,
         "init → prealign → refine →\nlandmarks → alignFull --tps",
         color=HAIRLINE, body_size=11)
    node(Inches(7.35), yB, Inches(2.4), h, "em_in_ccf.zarr", color=PURPLE)
    node(Inches(10.05), yB, Inches(2.7), h,
         "reference:\nuct_in_ccf.zarr",
         color=ATLAS, body_size=12, dim=True)

    arrow(s, Inches(2.75), yB + Emu(h.emu//2), Inches(3.05), yB + Emu(h.emu//2))
    arrow(s, Inches(7.05), yB + Emu(h.emu//2), Inches(7.35), yB + Emu(h.emu//2))
    arrow(s, Inches(10.05), yB + Emu(h.emu//2), Inches(9.75), yB + Emu(h.emu//2),
          color=ATLAS, weight=1.5)

    # Outcome strip at bottom
    out_top = Inches(6.15)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), out_top, Inches(12.2), Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = CARD_HI
    bar.line.color.rgb = HAIRLINE; bar.line.width = Pt(0.5)
    bar.shadow.inherit = False; bar.adjustments[0] = 0.5
    text(s, Inches(0.85), out_top, Inches(11.6), Inches(0.65),
         "Both volumes end up in CCF coordinates — no custom multi-stage abstraction.",
         size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    footer(s, n, total, kicker="Data flow")


def slide_invariants(p, n, total):
    s = add_slide(p)
    title(s, "Five invariants — true across modality and scale")
    subtitle(s, "If a feature breaks one of these, it is the wrong feature.")

    items = [
        ("01",  "State is the contract",
         "Every step reads and writes one state.json. Steps never talk to each other directly."),
        ("02",  "Reference is an abstraction",
         "All steps consume Reference{array, voxel_um, name}. BrainGlobe atlas and OME-Zarr are just two loaders."),
        ("03",  "Display strided, warp chunked",
         "Interactive napari uses ≤ 192³ preview. Production warp reads one chunk at a time. Full volumes never enter RAM."),
        ("04",  "Units stay in µm throughout",
         "EM at 4 nm is stored as 0.004 µm. No unit-tagging schema, no conversion layer."),
        ("05",  "One sample → one chain link → one state file",
         "Stage A and Stage B are two independent runs. The chain lives in --reference, not in the schema."),
    ]
    y = Inches(2.05); row_h = Inches(0.92)
    for num, head, body in items:
        # num chip
        c = card(s, Inches(0.55), y, Inches(0.85), row_h,
                 fill=CARD, border=ACCENT, border_w=1.0)
        text(s, Inches(0.55), y, Inches(0.85), row_h,
             num, size=22, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # body card
        card(s, Inches(1.55), y, Inches(11.2), row_h,
             fill=CARD, border=HAIRLINE)
        text(s, Inches(1.85), y + Inches(0.15), Inches(10.8), Inches(0.45),
             head, size=16, bold=True, color=INK)
        text(s, Inches(1.85), y + Inches(0.50), Inches(10.8), Inches(0.45),
             body, size=12, color=INK_DIM)
        y += row_h + Inches(0.08)

    footer(s, n, total, kicker="Design invariants")


def slide_reference(p, n, total):
    s = add_slide(p)
    title(s, "Reference abstraction — one interface, two loaders")
    subtitle(s, "Steps depend on a tiny dataclass. BrainGlobe atlas and OME-Zarr both populate it.")

    # Three columns: loaders | Reference dataclass | step consumers
    # Loaders (left)
    text(s, Inches(0.55), Inches(2.05), Inches(3.5), Inches(0.4),
         "LOADERS", size=11, bold=True, color=INK_DIM)
    c1 = card(s, Inches(0.55), Inches(2.5), Inches(3.5), Inches(1.4),
              border=ATLAS, border_w=1.0)
    text(s, Inches(0.75), Inches(2.65), Inches(3.1), Inches(0.4),
         "load_ccf(atlas_name)", size=14, bold=True, color=ATLAS, font="Consolas")
    text(s, Inches(0.75), Inches(3.05), Inches(3.1), Inches(0.4),
         "BrainGlobe atlas (current).", size=12, color=INK_DIM)
    text(s, Inches(0.75), Inches(3.40), Inches(3.1), Inches(0.4),
         "Used when state.atlas_name is set.", size=11, color=INK_DIM, italic=True)

    c2 = card(s, Inches(0.55), Inches(4.1), Inches(3.5), Inches(1.4),
              border=PURPLE, border_w=1.0)
    text(s, Inches(0.75), Inches(4.25), Inches(3.1), Inches(0.4),
         "load_zarr_reference(path)", size=14, bold=True, color=PURPLE, font="Consolas")
    text(s, Inches(0.75), Inches(4.65), Inches(3.1), Inches(0.4),
         "OME-Zarr volume (new).", size=12, color=INK_DIM)
    text(s, Inches(0.75), Inches(5.00), Inches(3.1), Inches(0.4),
         "Used when state.reference_zarr is set.", size=11, color=INK_DIM, italic=True)

    # Reference dataclass (middle)
    text(s, Inches(4.55), Inches(2.05), Inches(4.0), Inches(0.4),
         "INTERFACE", size=11, bold=True, color=INK_DIM)
    ref = card(s, Inches(4.55), Inches(2.5), Inches(4.0), Inches(3.0),
               fill=CARD_HI, border=ACCENT, border_w=1.5)
    text(s, Inches(4.55), Inches(2.65), Inches(4.0), Inches(0.4),
         "Reference", size=18, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
    # divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(5.05), Inches(3.15), Inches(3.0), Pt(1))
    div.fill.solid(); div.fill.fore_color.rgb = HAIRLINE
    div.line.fill.background(); div.shadow.inherit = False
    rich_text(s, Inches(4.85), Inches(3.30), Inches(3.4), Inches(2.0), [
        [("array",     {"font": "Consolas", "bold": True, "color": INK, "size": 13}),
         (": ndarray (z, y, x)", {"color": INK_DIM, "size": 12})],
        [("voxel_um",  {"font": "Consolas", "bold": True, "color": INK, "size": 13}),
         (": tuple[float, float, float]", {"color": INK_DIM, "size": 12})],
        [("name",      {"font": "Consolas", "bold": True, "color": INK, "size": 13}),
         (": str", {"color": INK_DIM, "size": 12})],
    ], line_spacing=1.5)

    # Step consumers (right)
    text(s, Inches(9.0), Inches(2.05), Inches(3.8), Inches(0.4),
         "CONSUMERS  (6 one-line edits)", size=11, bold=True, color=INK_DIM)
    steps = ["prealign", "refine", "landmarks", "alignFull", "export", "alignFull --tps"]
    for i, st in enumerate(steps):
        row_y = Inches(2.5 + i * 0.50)
        card(s, Inches(9.0), row_y, Inches(3.8), Inches(0.40),
             fill=CARD, border=HAIRLINE)
        text(s, Inches(9.2), row_y, Inches(3.6), Inches(0.40),
             st, size=12, color=INK, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)

    # Arrows: loaders → reference
    arrow(s, Inches(4.05), Inches(3.2), Inches(4.55), Inches(3.5), color=ATLAS)
    arrow(s, Inches(4.05), Inches(4.8), Inches(4.55), Inches(4.5), color=PURPLE)
    # Arrow: reference → consumers (single bus)
    arrow(s, Inches(8.55), Inches(4.0), Inches(9.0), Inches(4.0), color=ACCENT)

    # Bottom takeaway
    text(s, Inches(0.55), Inches(6.40), Inches(12.2), Inches(0.5),
         "Six 1-line edits in the step files. No step logic changes.",
         size=14, bold=True, color=GOOD, align=PP_ALIGN.CENTER)

    footer(s, n, total, kicker="Reference abstraction")


def slide_duality(p, n, total):
    s = add_slide(p)
    title(s, "Two paths through the data — display and warp")
    subtitle(s, "Both bounded in RAM. Different mechanisms because they have different jobs.")

    col_w = Inches(6.05); col_h = Inches(5.2); top = Inches(1.95)

    # DISPLAY card
    c1 = card(s, Inches(0.55), top, col_w, col_h,
              fill=CARD, border=ACCENT, border_w=1.0)
    text(s, Inches(0.85), top + Inches(0.25), col_w, Inches(0.5),
         "DISPLAY", size=12, bold=True, color=ACCENT)
    text(s, Inches(0.85), top + Inches(0.65), col_w, Inches(0.55),
         "prealign  ·  refine  ·  landmarks", size=18, bold=True, color=INK)
    rich_text(s, Inches(0.85), top + Inches(1.40), col_w - Inches(0.6), col_h - Inches(1.6), [
        [("Job.  ", {"bold": True, "color": ACCENT}),
         ("interactive picking — must be responsive.", {})],
        [("Loads.  ", {"bold": True, "color": ACCENT}),
         ("one strided preview, ≤ 192³ voxels.", {})],
        [("RAM.  ", {"bold": True, "color": ACCENT}),
         ("~30 MB regardless of source size.", {})],
        [("Trade-off.  ", {"bold": True, "color": ACCENT}),
         ("loses high-frequency detail.", {})],
        [("Mitigation.  ", {"bold": True, "color": ACCENT}),
         ("rigid + TPS captures the gross alignment; deformable refinement runs later on full resolution.", {})],
        [("Change needed.  ", {"bold": True, "color": WARN}),
         ("hardcoded 192³ → ", {}),
         ("--preview-size N",
          {"font": "Consolas", "color": INK}),
         (" so EM users can override.", {})],
    ], size=13, line_spacing=1.5)

    # WARP card
    c2 = card(s, Inches(6.75), top, col_w, col_h,
              fill=CARD, border=GOOD, border_w=1.0)
    text(s, Inches(7.05), top + Inches(0.25), col_w, Inches(0.5),
         "WARP", size=12, bold=True, color=GOOD)
    text(s, Inches(7.05), top + Inches(0.65), col_w, Inches(0.55),
         "alignFull (production)", size=18, bold=True, color=INK)
    rich_text(s, Inches(7.05), top + Inches(1.40), col_w - Inches(0.6), col_h - Inches(1.6), [
        [("Job.  ", {"bold": True, "color": GOOD}),
         ("produce the final aligned volume — must be exact.", {})],
        [("Loads.  ", {"bold": True, "color": GOOD}),
         ("one output chunk's source-bbox at a time.", {})],
        [("RAM.  ", {"bold": True, "color": GOOD}),
         ("bounded per worker  =  chunk × rotation overhead.", {})],
        [("Parallelism.  ", {"bold": True, "color": GOOD}),
         ("chunks dispatched via ThreadPoolExecutor.", {})],
        [("Already proven.  ", {"bold": True, "color": GOOD}),
         ("see align_full.py docstring — bypasses dask graph for TB inputs.", {})],
        [("Change needed.  ", {"bold": True, "color": WARN}),
         ("rigid-only → compose TPS (Step 3) and B-spline (Step 4).", {})],
    ], size=13, line_spacing=1.5)

    footer(s, n, total, kicker="Display vs. warp")


def slide_composition(p, n, total):
    s = add_slide(p)
    title(s, "Transform composition — the extensibility hinge")
    subtitle(s, "Each refinement step adds one component. warp_chunk walks the composition once per chunk.")

    # Pipeline boxes
    y = Inches(2.9); h = Inches(1.2)
    boxes = [
        ("Source\nvoxels",        Inches(0.55), Inches(2.0), PURPLE,  False),
        ("Rigid",                  Inches(2.95), Inches(2.0), ACCENT,  True),
        ("TPS\n(landmarks)",       Inches(5.35), Inches(2.0), GOOD,    True),
        ("B-spline\n(elastix / ANTs)", Inches(7.75), Inches(2.0), WARN, True),
        ("Reference\nvoxels",      Inches(10.30), Inches(2.4), ATLAS,  False),
    ]
    # rigid + tps + bspline get fatter boxes
    for label, x, w, color, is_op in boxes:
        bw = Inches(2.2) if is_op else Inches(2.4)
        c = card(s, x, y, bw, h, fill=CARD, border=color, border_w=1.25)
        text(s, x, y, bw, h, label,
             size=15, bold=True, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # status pill
        if is_op:
            done = label == "Rigid"
            stage = "DONE" if done else ("Step 3" if "TPS" in label else "Step 4")
            stage_color = GOOD if done else (ACCENT if stage == "Step 3" else WARN)
            chip(s, x + Inches(0.4), y + h + Inches(0.1),
                 bw - Inches(0.8), Inches(0.32), stage,
                 fill=CARD_HI, text_color=stage_color, size=10)

    # arrows
    yc = y + Emu(h.emu // 2)
    arrow(s, Inches(2.75), yc, Inches(2.95), yc)
    arrow(s, Inches(5.15), yc, Inches(5.35), yc)
    arrow(s, Inches(7.55), yc, Inches(7.75), yc)
    arrow(s, Inches(9.95), yc, Inches(10.30), yc, color=ATLAS, weight=1.5)

    # bottom explanation
    expl = card(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(2.05),
                fill=CARD, border=HAIRLINE)
    rich_text(s, Inches(0.85), Inches(5.00), Inches(11.7), Inches(1.85), [
        [("Interface.  ", {"bold": True, "color": ACCENT}),
         ("Each component answers one question: given a reference voxel coordinate, return the source voxel coordinate.", {"after": 8})],
        [("Per-chunk loop.  ", {"bold": True, "color": ACCENT}),
         ("warp_chunk applies the full composition once: chunk's reference grid → Rigid⁻¹ → TPS⁻¹ → B-spline⁻¹ → source sampling.", {"after": 8})],
        [("Independently optional.  ", {"bold": True, "color": ACCENT}),
         ("adding deformable does not refactor the rigid + TPS path. Each component is one method on a Transform class.", {})],
    ], size=13, line_spacing=1.4)

    footer(s, n, total, kicker="Transform composition")


def slide_dataflow(p, n, total):
    s = add_slide(p)
    title(s, "End-to-end data flow with state files")
    subtitle(s, "What lives on disk between commands.")

    col_w = Inches(6.05); col_h = Inches(5.30); top = Inches(1.85)

    # ---- Stage A ----
    left = Inches(0.55)
    card(s, left, top, col_w, col_h, fill=CARD, border=ACCENT, border_w=1.0)
    text(s, left + Inches(0.3), top + Inches(0.20), col_w, Inches(0.4),
         "STAGE A", size=12, bold=True, color=ACCENT)
    text(s, left + Inches(0.3), top + Inches(0.55), col_w, Inches(0.5),
         "µCT → Allen CCF", size=18, bold=True, color=INK)
    rich_text(s, left + Inches(0.3), top + Inches(1.20),
              col_w - Inches(0.6), col_h - Inches(1.35), [
        [("1.  vol2atlas init",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11}),
         ("  sample_uct.zarr  -s state_uct.json",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        writes   state_uct.json   ",
          {"font": "Consolas", "color": INK_DIM, "size": 9}),
         ("{sample_zarr, voxel_um, atlas_name}",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9, "after": 10})],
        [("2.  vol2atlas prealign / refine / landmarks",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11}),
         ("  state_uct.json",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        updates  state_uct.json   ",
          {"font": "Consolas", "color": INK_DIM, "size": 9}),
         ("{transform, landmarks, ccf_crop_bbox}",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9, "after": 10})],
        [("3.  vol2atlas alignFull",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11}),
         ("  state_uct.json  -o out/uct_in_ccf.zarr",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        writes   uct_in_ccf.zarr   ",
          {"font": "Consolas", "color": ATLAS, "size": 9}),
         ("(chunked, CCF voxel grid)",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9})],
    ], size=11, line_spacing=1.35)

    # ---- Stage B ----
    left = Inches(6.75)
    card(s, left, top, col_w, col_h, fill=CARD, border=PURPLE, border_w=1.0)
    text(s, left + Inches(0.3), top + Inches(0.20), col_w, Inches(0.4),
         "STAGE B", size=12, bold=True, color=PURPLE)
    text(s, left + Inches(0.3), top + Inches(0.55), col_w, Inches(0.5),
         "EM → µCT (already in CCF)", size=18, bold=True, color=INK)
    rich_text(s, left + Inches(0.3), top + Inches(1.20),
              col_w - Inches(0.6), col_h - Inches(1.35), [
        [("1.  vol2atlas init  sample_em.zarr  -s state_em.json  \\",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11})],
        [("           --reference out/uct_in_ccf.zarr  --voxel-um 0.004",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        writes   state_em.json   ",
          {"font": "Consolas", "color": INK_DIM, "size": 9}),
         ("{sample_zarr, voxel_um, reference_zarr}",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9, "after": 10})],
        [("2.  vol2atlas prealign / refine / landmarks",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11}),
         ("  state_em.json",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        updates  state_em.json   ",
          {"font": "Consolas", "color": INK_DIM, "size": 9}),
         ("{transform, landmarks}",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9, "after": 10})],
        [("3.  vol2atlas alignFull --tps",
          {"font": "Consolas", "color": INK, "bold": True, "size": 11}),
         ("  state_em.json  -o out/em_in_ccf.zarr",
          {"font": "Consolas", "color": INK, "size": 11})],
        [("        writes   em_in_ccf.zarr   ",
          {"font": "Consolas", "color": PURPLE, "size": 9}),
         ("(chunked, CCF voxel grid)",
          {"font": "Consolas", "color": INK_DIM, "italic": True, "size": 9})],
    ], size=11, line_spacing=1.35)

    footer(s, n, total, kicker="CLI flow")


def slide_anti_scope(p, n, total):
    s = add_slide(p)
    title(s, "What this strategy explicitly does NOT do")
    subtitle(s, "Anti-scope. Naming each is as important as the things we do.")

    items = [
        ("No unit-tagging in state",
         "voxel_unit: 'nm'|'um' is a footgun that propagates checks everywhere. Float µm is sufficient."),
        ("No hardcoded UI thresholds",
         "Slider step, preview size, landmark dot size all derive from sample_voxel_um. Remove 50.0 and 192³ constants."),
        ("No deformable registration on TB volumes",
         "elastix / ANTs run on a downsampled GB-level. The B-spline result composes into the chunkwise warp."),
        ("No nm-precision landmark UI",
         "Landmarks pin the rigid + coarse warp. nm precision comes from deformable refinement, not the picker."),
        ("No intra-stack EM section assembly",
         "Input must already be a 3D OME-Zarr. Section assembly is a different problem — TrakEM2 / ASAP / render-python."),
        ("No two-stage abstraction in state.json",
         "Two stages = two independent runs. Chaining lives in --reference, not in the schema."),
        ("No brainreg path for EM",
         "brainreg requires a BrainGlobe atlas name. The brainreg subcommand errors cleanly when reference_zarr is set."),
    ]
    y = Inches(2.0); row_h = Inches(0.66)
    for head, body in items:
        card(s, Inches(0.55), y, Inches(12.2), row_h,
             fill=CARD, border=HAIRLINE)
        # X badge
        text(s, Inches(0.85), y, Inches(0.4), row_h,
             "✗", size=18, bold=True, color=WARN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # head
        text(s, Inches(1.35), y + Inches(0.10), Inches(4.8), Inches(0.45),
             head, size=13, bold=True, color=INK)
        text(s, Inches(1.35), y + Inches(0.35), Inches(4.8), Inches(0.30),
             "", size=8, color=INK_DIM)   # spacer
        # body
        text(s, Inches(6.4), y + Inches(0.10), Inches(6.2), row_h,
             body, size=11, color=INK_DIM, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + Inches(0.06)

    footer(s, n, total, kicker="Anti-scope")


def slide_roadmap(p, n, total):
    s = add_slide(p)
    title(s, "Implementation roadmap")
    subtitle(s, "Four steps. Each independently shippable and usable.")

    # Horizontal timeline with 4 milestones
    track_y = Inches(3.8)
    # track line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1.0), track_y, Inches(11.3), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = HAIRLINE
    line.line.fill.background(); line.shadow.inherit = False

    steps = [
        ("01", "Reference\nabstraction",
         "load_reference + Reference dataclass.\nSix 1-line edits across step files.",
         "Enables Stage B with rigid at any scale.",
         GOOD, "Smallest change · highest leverage"),
        ("02", "Voxel-size-\naware UI",
         "Sliders derive step from sample_voxel_um.\n--preview-size flag for prealign/refine/landmarks.",
         "Makes the UI usable on EM. No schema impact.",
         GOOD, "Pure UI change"),
        ("03", "alignFull --tps",
         "Extend _warp_one_chunk to compose landmark\nTPS displacement per chunk. ~50 LOC.",
         "Production EM path at TB scale.",
         ACCENT, "TPS parity with export"),
        ("04", "alignFull\n--bspline",
         "Compose elastix / ANTs B-spline output into\nthe chunkwise warp.",
         "Unlocks nm-precision deformable refinement.",
         WARN, "Separate PR — depends on 1-3"),
    ]
    col_x = [Inches(0.65), Inches(3.65), Inches(6.65), Inches(9.65)]
    col_w = Inches(3.0)
    for (num, name, body, value, color, kicker), x in zip(steps, col_x):
        # node dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
            x + Inches(1.35), track_y - Inches(0.13), Inches(0.32), Inches(0.32))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.color.rgb = BG; dot.line.width = Pt(2.5)
        dot.shadow.inherit = False
        # number above the dot
        text(s, x, Inches(1.95), col_w, Inches(0.4),
             num, size=14, bold=True, color=color,
             align=PP_ALIGN.CENTER)
        # title card above track
        card(s, x, Inches(2.35), col_w, Inches(1.25),
             fill=CARD, border=color, border_w=1.0)
        text(s, x, Inches(2.5), col_w, Inches(1.0),
             name, size=15, bold=True, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # body card below track
        card(s, x, Inches(4.20), col_w, Inches(2.3),
             fill=CARD, border=HAIRLINE)
        text(s, x + Inches(0.2), Inches(4.32), col_w - Inches(0.4), Inches(1.0),
             body, size=11, color=INK_DIM, line_spacing=1.4)
        text(s, x + Inches(0.2), Inches(5.40), col_w - Inches(0.4), Inches(0.5),
             value, size=11, bold=True, color=INK, line_spacing=1.3)
        text(s, x + Inches(0.2), Inches(5.90), col_w - Inches(0.4), Inches(0.45),
             "→ " + kicker, size=10, bold=True, color=color, italic=True)

    footer(s, n, total, kicker="Roadmap")


def slide_summary(p, n, total):
    s = add_slide(p)
    title(s, "Summary  ·  one pipeline, every scale")
    subtitle(s, "Three architectural moves unlock EM at TB scale.")

    rows = [
        ("Generalize the reference",
         "BrainGlobe atlas and OME-Zarr behind one Reference interface. Six 1-line edits.",
         GOOD),
        ("Derive UI from voxel size",
         "Slider step, preview stride, landmark size all scale with sample_voxel_um.",
         GOOD),
        ("Compose the warp",
         "alignFull walks Rigid ∘ TPS ∘ B-spline per chunk. RAM bounded; no whole-volume load.",
         ACCENT),
    ]
    y = Inches(2.15); h = Inches(1.3)
    for head, body, color in rows:
        card(s, Inches(0.55), y, Inches(12.2), h,
             fill=CARD, border=color, border_w=1.0)
        # left rail
        rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.55), y, Inches(0.10), h)
        rail.fill.solid(); rail.fill.fore_color.rgb = color
        rail.line.fill.background(); rail.shadow.inherit = False
        text(s, Inches(0.85), y + Inches(0.20), Inches(11.5), Inches(0.5),
             head, size=18, bold=True, color=INK)
        text(s, Inches(0.85), y + Inches(0.70), Inches(11.5), Inches(0.5),
             body, size=13, color=INK_DIM)
        y += h + Inches(0.20)

    # outcome strip
    out = card(s, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.55),
               fill=CARD_HI, border=ACCENT, border_w=1.0)
    text(s, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.55),
         "EM lands in CCF coordinates through the same CLI as µCT, with no TB-scale RAM cost.",
         size=13, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total, kicker="Summary")


# ============================================================================
def main():
    p = new_pres()
    builders = [
        slide_title,
        slide_problem,
        slide_chain,
        slide_invariants,
        slide_reference,
        slide_duality,
        slide_composition,
        slide_dataflow,
        slide_anti_scope,
        slide_roadmap,
        slide_summary,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        b(p, i, total)

    out = Path(__file__).resolve().parent.parent / "vol2atlas_strategy.pptx"
    p.save(out)
    print(f"wrote {out}  ({total} slides)")


if __name__ == "__main__":
    main()
