"""Generate vol2atlas_method.pptx — 2 slides (µCT→CCF, EM→aligned µCT).

Same dark theme as make_strategy_slides.py.
Run:  python scripts/make_method_slides.py
Out:  vol2atlas_method.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

BG       = RGBColor(0x10, 0x10, 0x14)
CARD     = RGBColor(0x1C, 0x1C, 0x22)
CARD_HI  = RGBColor(0x24, 0x24, 0x2C)
INK      = RGBColor(0xEE, 0xEE, 0xEE)
INK_DIM  = RGBColor(0x9A, 0x9A, 0xA6)
ACCENT   = RGBColor(0x8A, 0xB4, 0xFF)
GOOD     = RGBColor(0x81, 0xC9, 0x95)
WARN     = RGBColor(0xF2, 0x8B, 0x82)
ATLAS    = RGBColor(0xFD, 0xD6, 0x63)
PURPLE   = RGBColor(0xC5, 0x8A, 0xF9)
HAIRLINE = RGBColor(0x33, 0x33, 0x3D)
FONT     = "Inter"

SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)


def add_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def text(slide, left, top, width, height, body, *,
         size=13, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align; para.line_spacing = line_spacing
        r = para.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def rich_text(slide, left, top, width, height, paragraphs,
              size=13, color=INK, line_spacing=1.35, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for pi, runs in enumerate(paragraphs):
        para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = line_spacing
        for txt, kw in runs:
            r = para.add_run(); r.text = txt
            r.font.name  = kw.get("font",  FONT)
            r.font.size  = Pt(kw.get("size",  size))
            r.font.bold  = kw.get("bold",  False)
            r.font.italic = kw.get("italic", False)
            r.font.color.rgb = kw.get("color", color)
        para.space_after = Pt(kw.get("after", 4))
    return tb


def card(slide, left, top, width, height, *,
         fill=CARD, border=HAIRLINE, border_w=0.75, rounded=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    if rounded:
        sh.adjustments[0] = 0.06
    return sh


def chip(slide, left, top, width, height, label, *,
         fill=CARD_HI, text_color=ACCENT, size=10, bold=True):
    sh = card(slide, left, top, width, height, fill=fill, border=fill, rounded=True)
    sh.adjustments[0] = 0.5
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p_ = tf.paragraphs[0]; p_.alignment = PP_ALIGN.CENTER
    r = p_.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = text_color


def arrow(slide, x1, y1, x2, y2, *, color=INK_DIM, weight=1.25, head=True):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color; line.line.width = Pt(weight)
    if head:
        from lxml import etree
        from pptx.oxml.ns import qn
        ln = line.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('h', 'med')


def title(slide, body):
    text(slide, Inches(0.55), Inches(0.4), Inches(12.3), Inches(0.85),
         body, size=30, bold=True, color=INK)


def subtitle(slide, body, color=INK_DIM):
    text(slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(0.5),
         body, size=14, color=color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(1.62), Inches(0.65), Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False


def footer(slide, n, total, kicker=""):
    text(slide, Inches(0.55), Inches(7.13), Inches(8), Inches(0.3),
         kicker, size=9, color=INK_DIM)
    text(slide, Inches(11.5), Inches(7.13), Inches(1.4), Inches(0.3),
         f"{n} / {total}", size=10, color=INK_DIM, align=PP_ALIGN.RIGHT)


def pipeline_strip(s, y, h, nodes, accent):
    """Draw a horizontal pipeline of node cards with arrows between them."""
    x = Inches(0.55)
    total_w = Inches(12.2)
    n = len(nodes)
    gap = Inches(0.2)
    node_w = Emu((total_w.emu - gap.emu * (n - 1)) // n)
    for i, (label, color) in enumerate(nodes):
        xi = Emu(x.emu + i * (node_w.emu + gap.emu))
        c = card(s, xi, y, node_w, h, fill=CARD, border=color, border_w=1.25)
        text(s, xi, y, node_w, h, label, size=13, bold=True, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ax1 = Emu(xi.emu + node_w.emu)
            ax2 = Emu(ax1.emu + gap.emu)
            arrow(s, ax1, Emu(y.emu + h.emu // 2),
                  ax2, Emu(y.emu + h.emu // 2), color=accent, weight=1.5)


# ============================================================================
# SLIDE 1 — µCT → CCF
# ============================================================================
def slide_uct(p, n, total):
    s = add_slide(p)
    title(s, "Stage A  ·  µCT  →  Allen CCF")
    subtitle(s, "Working today — tested on real µCT. Rigid + optional landmark TPS is solid; the deformable refiners are wrappers with known limits.",
             color=INK_DIM)

    # Top kicker
    text(s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.4),
         "INPUT  µCT OME-Zarr (µm voxels, GB scale)   →   "
         "REFERENCE  BrainGlobe Allen CCF   →   "
         "OUTPUT  warped sample on the CCF voxel grid",
         size=11, bold=True, color=ACCENT)

    # Pipeline strip
    pipeline_strip(s, Inches(2.35), Inches(0.85), [
        ("init",       ATLAS),
        ("prealign",   ACCENT),
        ("refine",     ACCENT),
        ("landmarks",  ACCENT),
        ("alignFull",  GOOD),
    ], accent=INK_DIM)

    # Two-column body
    col_top = Inches(3.30); col_h = Inches(3.70)
    col_w   = Inches(6.05)

    # LEFT — method
    card(s, Inches(0.55), col_top, col_w, col_h, fill=CARD, border=ACCENT, border_w=1.0)
    text(s, Inches(0.85), col_top + Inches(0.20), col_w - Inches(0.5), Inches(0.4),
         "METHOD", size=11, bold=True, color=ACCENT)
    rich_text(s, Inches(0.85), col_top + Inches(0.60),
              col_w - Inches(0.6), col_h - Inches(0.75), [
        [("Reference.  ", {"bold": True, "color": ACCENT}),
         ("Allen CCF via BrainGlobe (atlas_name).", {})],
        [("Display.  ", {"bold": True, "color": ACCENT}),
         ("napari MIP, strided preview ≤ 192³ — RAM bounded regardless of source size.", {})],
        [("Rigid pose.  ", {"bold": True, "color": ACCENT}),
         ("6 sliders (3 translations µm, 3 rotations °) + axis-flip checkboxes; live resample onto CCF grid.", {})],
        [("Crop ROI.  ", {"bold": True, "color": ACCENT}),
         ("6 min/max sliders crop the CCF to a bbox around the sample.", {})],
        [("Refine.  ", {"bold": True, "color": ACCENT}),
         ("±3 mm / ±30° in axial / coronal / sagittal ortho views.", {})],
        [("Landmarks.  ", {"bold": True, "color": ACCENT}),
         ("optional Procrustes / Kabsch refit on ≥3 paired clicks.", {})],
        [("Production warp.  ", {"bold": True, "color": ACCENT}),
         ("alignFull warps the full OME-Zarr at every pyramid level, one output chunk at a time via dask_image.ndinterp.", {})],
    ], size=11, line_spacing=1.25)

    # RIGHT — math & data
    card(s, Inches(6.75), col_top, col_w, col_h, fill=CARD, border=GOOD, border_w=1.0)
    text(s, Inches(7.05), col_top + Inches(0.20), col_w - Inches(0.5), Inches(0.4),
         "STATE & MATH", size=11, bold=True, color=GOOD)
    rich_text(s, Inches(7.05), col_top + Inches(0.60),
              col_w - Inches(0.6), col_h - Inches(0.75), [
        [("Coords.  ", {"bold": True, "color": GOOD}),
         ("everything in (z, y, x) µm. No unit-tagging schema.", {})],
        [("Rigid.  ", {"bold": True, "color": GOOD}),
         ("output_µm = R · F · (input_µm − c) + c + t,  with R = intrinsic ZYX Euler, F = ±1 axis flips, c = sample center, t = translation.", {})],
        [("Optional TPS.  ", {"bold": True, "color": GOOD}),
         ("thin-plate spline on the residuals between rigid prediction and landmark clicks; vanishes far from landmarks → rigid preserved.", {})],
        [("State file.  ", {"bold": True, "color": GOOD}),
         ("one state.json per stage: {sample_zarr, voxel_um, sample_level, atlas_name, transform, landmarks, ccf_crop_bbox, history}. Resume any step.", {})],
        [("Output.  ", {"bold": True, "color": GOOD}),
         ("uct_in_ccf.zarr — multiscale OME-Zarr, chunked, CCF voxel grid. Used directly as Stage B's reference.", {})],
    ], size=11, line_spacing=1.25)

    footer(s, n, total, kicker="Stage A — µCT to CCF")


# ============================================================================
# SLIDE 2 — EM → aligned µCT
# ============================================================================
def slide_em(p, n, total):
    s = add_slide(p)
    title(s, "Stage B  ·  EM  →  aligned µCT  (proposed)")
    subtitle(s, "Roadmap — not yet implemented and not yet tested on EM data.",
             color=WARN)

    # ROADMAP banner
    banner = card(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.45),
                  fill=CARD_HI, border=WARN, border_w=1.0)
    text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.45),
         "⚠  ROADMAP  —  --reference, alignFull --tps, voxel-aware UI all unmerged. "
         "No EM volume has been run through this pipeline.",
         size=11, bold=True, color=WARN,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Top kicker (moved down)
    text(s, Inches(0.55), Inches(2.40), Inches(12), Inches(0.4),
         "PROPOSED INPUT  EM OME-Zarr (nm voxels, TB scale)   →   "
         "PROPOSED REFERENCE  uct_in_ccf.zarr   →   "
         "PROPOSED OUTPUT  warped EM on the CCF voxel grid",
         size=11, bold=True, color=PURPLE)

    # Pipeline strip
    pipeline_strip(s, Inches(2.85), Inches(0.75), [
        ("init --reference",  PURPLE),
        ("prealign",          ACCENT),
        ("refine",            ACCENT),
        ("landmarks",         ACCENT),
        ("alignFull --tps",   GOOD),
    ], accent=PURPLE)

    # Two-column body
    col_top = Inches(3.75); col_h = Inches(3.25)
    col_w   = Inches(6.05)

    # LEFT — what would change vs Stage A
    card(s, Inches(0.55), col_top, col_w, col_h, fill=CARD, border=PURPLE, border_w=1.0)
    text(s, Inches(0.85), col_top + Inches(0.20), col_w - Inches(0.5), Inches(0.4),
         "WHAT WOULD CHANGE vs STAGE A  (to be built)", size=11, bold=True, color=PURPLE)
    rich_text(s, Inches(0.85), col_top + Inches(0.60),
              col_w - Inches(0.6), col_h - Inches(0.75), [
        [("Reference loader.  ", {"bold": True, "color": PURPLE}),
         ("add load_zarr_reference(path) alongside load_ccf. ~6 one-line edits.", {})],
        [("Voxel scale.  ", {"bold": True, "color": PURPLE}),
         ("4 nm stored as 0.004 µm. No unit conversion.", {})],
        [("UI scaling.  ", {"bold": True, "color": PURPLE}),
         ("slider step, preview stride, landmark marker size derive from sample_voxel_um.", {})],
        [("Warp composition.  ", {"bold": True, "color": PURPLE}),
         ("alignFull --tps composes Rigid ∘ TPS per chunk. B-spline plugs in later.", {})],
    ], size=11, line_spacing=1.3)

    # RIGHT — design target (carried over from Stage A's proven engine)
    card(s, Inches(6.75), col_top, col_w, col_h, fill=CARD, border=GOOD, border_w=1.0)
    text(s, Inches(7.05), col_top + Inches(0.20), col_w - Inches(0.5), Inches(0.4),
         "DESIGN TARGET  (reuses Stage A engine — proven on µCT)", size=11, bold=True, color=GOOD)
    rich_text(s, Inches(7.05), col_top + Inches(0.60),
              col_w - Inches(0.6), col_h - Inches(0.75), [
        [("RAM bound (proven on µCT).  ", {"bold": True, "color": GOOD}),
         ("warp reads one output chunk's source-bbox at a time.", {})],
        [("Parallelism (proven on µCT).  ", {"bold": True, "color": GOOD}),
         ("chunks via ThreadPoolExecutor; bypasses dask graph build.", {})],
        [("State (designed).  ", {"bold": True, "color": GOOD}),
         ("state_em.json adds reference_zarr. state_uct.json unchanged.", {})],
        [("Open risks.  ", {"bold": True, "color": WARN}),
         ("nm-scale slider/click ergonomics, OME-Zarr level-selection on TB inputs, "
          "TPS conditioning on nm landmark spreads — all untested.", {})],
        [("Out of scope.  ", {"bold": True, "color": WARN}),
         ("EM section assembly (TrakEM2 / render-python).", {})],
    ], size=11, line_spacing=1.3)

    footer(s, n, total, kicker="Stage B — EM to aligned µCT")


# ============================================================================
def main():
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H

    slide_uct(p, 1, 2)
    slide_em (p, 2, 2)

    out = Path(__file__).resolve().parent.parent / "vol2atlas_method.pptx"
    p.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
