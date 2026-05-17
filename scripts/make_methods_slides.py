"""Generates methods_survey.pptx — overview of the registration methods
in the literature relevant for mouse-brain → Allen CCF, with a focus
on partial-sample applicability.

Run:
    python scripts/make_methods_slides.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent.parent / "methods_survey.pptx"

# ---- theme ---------------------------------------------------------------
BG       = RGBColor(0x10, 0x10, 0x14)
PANEL    = RGBColor(0x1c, 0x1c, 0x22)
FG       = RGBColor(0xee, 0xee, 0xee)
MUTED    = RGBColor(0x9a, 0x9a, 0xa6)
ACCENT   = RGBColor(0x8a, 0xb4, 0xff)   # blue
OK       = RGBColor(0x9c, 0xe8, 0xa0)   # green
WARN     = RGBColor(0xff, 0x9c, 0x8a)   # red-orange
YELLOW   = RGBColor(0xff, 0xd2, 0x6c)
CODE_BG  = RGBColor(0x16, 0x18, 0x1d)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_bg(s):
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return bg


def add_text(s, x, y, w, h, text, *, size=18, bold=False,
             color=FG, align=PP_ALIGN.LEFT, mono=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        f = r.font; f.size = Pt(size); f.bold = bold
        f.color.rgb = color
        f.name = "JetBrains Mono" if mono else "Inter"
    return tb


def add_panel(s, x, y, w, h, *, fill=PANEL, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.05
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    return sh


def add_title(s, title, subtitle=None):
    add_text(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             title, size=28, bold=True)
    if subtitle:
        add_text(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=MUTED)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(1.32),
                              Inches(0.6), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def bullets(s, x, y, w, h, items, *, size=14, lead=ACCENT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run(); r1.text = "•  "
            r1.font.size = Pt(size); r1.font.color.rgb = lead
            r2 = p.add_run(); r2.text = head
            r2.font.size = Pt(size); r2.font.bold = True
            r2.font.color.rgb = FG; r2.font.name = "Inter"
            r3 = p.add_run(); r3.text = "  " + body
            r3.font.size = Pt(size); r3.font.color.rgb = FG; r3.font.name = "Inter"
        else:
            r1 = p.add_run(); r1.text = "•  "
            r1.font.size = Pt(size); r1.font.color.rgb = lead
            r2 = p.add_run(); r2.text = item
            r2.font.size = Pt(size); r2.font.color.rgb = FG; r2.font.name = "Inter"


def method_card(s, x, y, w, h, *, name, cite, body, partial,
                repo=None, border=ACCENT):
    add_panel(s, x, y, w, h, border=border)
    add_text(s, x + Inches(0.2), y + Inches(0.12),
             w - Inches(0.4), Inches(0.4),
             name, size=15, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.5),
             w - Inches(0.4), Inches(0.3),
             cite, size=10, color=MUTED, mono=True)
    add_text(s, x + Inches(0.2), y + Inches(0.85),
             w - Inches(0.4), h - Inches(1.4),
             body, size=11, color=FG)
    badge = "PARTIAL SAMPLE: ✓" if partial else "PARTIAL SAMPLE: ✗"
    badge_color = OK if partial else WARN
    add_text(s, x + Inches(0.2), y + h - Inches(0.45),
             w - Inches(0.4), Inches(0.3),
             badge, size=10, bold=True, color=badge_color, mono=True)
    if repo:
        add_text(s, x + Inches(0.2), y + h - Inches(0.25),
                 w - Inches(0.4), Inches(0.2),
                 repo, size=8, color=MUTED, mono=True)


# ==========================================================================
# 1. TITLE
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.1),
         "Mouse-brain → CCF registration", size=44, bold=True)
add_text(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
         "Methods landscape", size=22, color=ACCENT)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
         "What exists in the literature, what works on partial-hemisphere µCT, "
         "and how each one fits into the vol2atlas pipeline.",
         size=14, color=MUTED)
add_text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
         "Alberto Mittone — APS / Argonne National Laboratory",
         size=12, color=MUTED)
add_text(s, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.4),
         "2026-05-17", size=12, color=MUTED)

# ==========================================================================
# 2. PROBLEM
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "The task & why it is non-trivial",
          "Map a µCT mouse brain volume into the Allen CCFv3 atlas grid "
          "so downstream tools can use region labels.")
bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5), [
    ("Input.",           "Multi-TB OME-Zarr µCT volume (~2.74 µm/voxel at level 2). "
                          "Often only one hemisphere or a partial brain."),
    ("Reference.",       "Allen Mouse Brain CCFv3, 25 µm/voxel, full brain (528×320×456)."),
    ("Hard parts.",      "(i) cross-modality intensities (µCT looks nothing like the LSFM-derived atlas); "
                          "(ii) partial samples (most tools assume a full brain or hemisphere); "
                          "(iii) huge data volumes (anything that loads the full level into RAM is dead)."),
    ("Method families.", "intensity-driven deformable · landmark/spline-driven · deep learning · "
                          "atlas-specific pipelines · interactive GUIs · slice-only tools."),
    ("Practical question.",
                         "Which one (or combination) gives the best CCF alignment for "
                          "a partial-hemisphere µCT sample, in reasonable time?"),
])

# ==========================================================================
# 3. INTENSITY-DRIVEN DEFORMABLE
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Intensity-driven deformable registration",
          "General-purpose tools, not mouse-brain-specific. Maximise an "
          "intensity similarity metric (CC, MI, MSE) over a deformation field.")

w, h, gap = Inches(3.05), Inches(2.55), Inches(0.13)
cards = [
    dict(name="ANTs SyN / SyNOnly",
         cite="Avants et al. 2008 (Med Image Anal)",
         body="Diffeomorphic symmetric normalization. SyN does rigid+affine+SyN; "
              "SyNOnly skips rigid+affine and is the right choice when you have "
              "a good prealignment. Supports masks on both fixed/moving.",
         partial=True,
         repo="pip install antspyx · github.com/ANTsX/ANTs"),
    dict(name="elastix / SimpleElastix",
         cite="Klein et al. 2010 (IEEE TMI)",
         body="B-spline FFD with many configurable optimizers + metrics. "
              "Wrapped in SimpleElastix (Python). Used by ClearMap2 and "
              "brainregister downstream. Mask support is explicit.",
         partial=True,
         repo="pip install SimpleITK · elastix.dev"),
    dict(name="NiftyReg (f3d, aladin)",
         cite="Modat et al. 2010 (CMPB)",
         body="CUDA-accelerated B-spline registration; engine inside BrainGlobe's "
              "brainreg. Fast, but the brainreg wrapper assumes a full brain.",
         partial=False,
         repo="github.com/KCL-BMEIS/niftyreg"),
    dict(name="Greedy",
         cite="Yushkevich, in ITK-SNAP",
         body="Greedy diffeomorphic, designed as a fast ANTs alternative for "
              "clinical workflows. Supports masks. Less documentation than ANTs.",
         partial=True,
         repo="github.com/pyushkevich/greedy"),
]
x0, y0 = Inches(0.4), Inches(1.5)
for i, c in enumerate(cards):
    col = i % 4
    method_card(s, x0 + col * (w + gap), y0, w, h, **c)

# next row: brief context note
add_panel(s, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.6),
          fill=PANEL, border=MUTED)
add_text(s, Inches(0.6), Inches(4.45), Inches(12.0), Inches(0.4),
         "Notes for partial-hemisphere data",
         size=15, bold=True, color=ACCENT)
bullets(s, Inches(0.6), Inches(4.9), Inches(12.0), Inches(2.0), [
    "Always pass brain masks — without them, the metric scores the empty "
    "(zero) regions around the sample as 'missing tissue' and the deformation "
    "field stretches into them. ANTs auto-masks via `ants.get_mask`.",
    "Never combine a deformable stage with an affine that re-fits the global pose "
    "(SyN, ANTs `Affine`, NiftyReg `aladin`) — it will overwrite your prealignment. "
    "Use SyNOnly / elastix-BSpline-only.",
    "Heavy regularization helps: bigger `flow_sigma`/`total_sigma` in ANTs, "
    "lower `FinalBSplineInterpolationOrder` in elastix, fewer iterations "
    "in coarse-to-fine schedules.",
], size=12)

# ==========================================================================
# 4. DEEP-LEARNING REGISTRATION
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Deep-learning registration",
          "CNN regresses the deformation field directly. Fast at inference; "
          "training data is the catch.")

cards = [
    dict(name="VoxelMorph",
         cite="Balakrishnan et al. 2019 (IEEE TMI)",
         body="Unsupervised CNN, learns a probabilistic diffeomorphic warp. "
              "Many forks; pretrained weights exist for human MRI, less so for "
              "mouse brain.",
         partial=False,
         repo="github.com/voxelmorph/voxelmorph"),
    dict(name="SynthMorph",
         cite="Hoffmann et al. 2022 (IEEE TMI)",
         body="VoxelMorph trained on synthetic images so it is contrast-invariant. "
              "More robust to µCT-vs-LSFM intensity mismatch. Still needs full-volume "
              "training data to fine-tune on partial hemispheres.",
         partial=False,
         repo="github.com/voxelmorph/voxelmorph (synthmorph branch)"),
    dict(name="mBrainAligner 3D U-Net",
         cite="Qu et al. 2022 (Nat Methods)",
         body="DL component used INSIDE mBrainAligner to predict landmark positions "
              "for the CLM (Coherent Landmark Mapping) step. Not a standalone "
              "registration tool, but the network drives the half-brain pipeline.",
         partial=True,
         repo="github.com/Vaa3D/vaa3d_tools/.../mBrainAligner/src/3D_U-Net"),
    dict(name="DeepReg / TransMorph",
         cite="Various 2021–2024",
         body="Transformer- or attention-based registration networks. Mostly "
              "validated on human MRI; mouse-brain checkpoints are rare. Worth "
              "watching but not turnkey.",
         partial=False,
         repo="github.com/junyuchen245/TransMorph"),
]
x0, y0 = Inches(0.4), Inches(1.5)
for i, c in enumerate(cards):
    col = i % 4
    method_card(s, x0 + col * (w + gap), y0, w, h, **c)

add_panel(s, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.6),
          fill=PANEL, border=MUTED)
add_text(s, Inches(0.6), Inches(4.45), Inches(12.0), Inches(0.4),
         "Bottom line for your case",
         size=15, bold=True, color=ACCENT)
bullets(s, Inches(0.6), Inches(4.9), Inches(12.0), Inches(2.0), [
    "No pretrained network for µCT mouse brain → CCF exists publicly today "
    "(2026-05). Training one needs labelled pairs you do not have.",
    "If you ever generate ~100 aligned µCT mouse brains, SynthMorph or "
    "TransMorph would be the obvious targets to fine-tune.",
    "Until then, intensity-driven (ANTs / elastix) is more practical than DL.",
], size=12)

# ==========================================================================
# 5. ATLAS-SPECIFIC PIPELINES
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Mouse-brain-specific atlas pipelines",
          "Tools that bundle a CCF/atlas + registration engine + sample-data "
          "preprocessing into one workflow.")

cards = [
    dict(name="mBrainAligner",
         cite="Qu et al. 2022 (Nat Methods)",
         body="Two-stage (global + local) with optional landmark refinement. "
              "Ships explicit half-brain target (`examples/target/half/`) and "
              "modality-specific configs. The only published tool that targets "
              "partial-hemisphere mouse brain → CCF directly.",
         partial=True,
         repo="github.com/Vaa3D/vaa3d_tools/.../mBrainAligner"),
    dict(name="BrainGlobe brainreg",
         cite="Tyson et al. 2022 (Sci Reports)",
         body="Wraps NiftyReg with the BrainGlobe atlas API. CLI-driven, easy "
              "to install, good defaults — but `--brain_geometry` only knows "
              "`full / hemisphere_l / hemisphere_r`, so partial samples distort.",
         partial=False,
         repo="pip install brainreg · brainglobe.info"),
    dict(name="brainregister",
         cite="Steven J. West et al. — int-brain-lab",
         body="Wraps elastix instead of NiftyReg, otherwise similar to brainreg. "
              "Same full-brain assumption baked into the pre-shipped parameter "
              "files; you would need custom param-maps for partial samples.",
         partial=False,
         repo="github.com/int-brain-lab/brainregister"),
    dict(name="ClearMap2",
         cite="Renier et al. 2016 / Kirst et al. 2020",
         body="Whole-pipeline for cleared brain (LSFM): preprocessing + elastix "
              "+ region quantification. Assumes whole brain. Useful as a reference "
              "for elastix parameter files even when not using the rest.",
         partial=False,
         repo="github.com/ChristophKirst/ClearMap2"),
]
x0, y0 = Inches(0.4), Inches(1.5)
for i, c in enumerate(cards):
    col = i % 4
    method_card(s, x0 + col * (w + gap), y0, w, h, **c)

add_panel(s, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.6),
          fill=PANEL, border=YELLOW)
add_text(s, Inches(0.6), Inches(4.45), Inches(12.0), Inches(0.4),
         "Key takeaway",
         size=15, bold=True, color=YELLOW)
bullets(s, Inches(0.6), Inches(4.9), Inches(12.0), Inches(2.0), [
    "mBrainAligner is the ONLY off-the-shelf pipeline with explicit "
    "partial-hemisphere support out of the box (the `half.7z` target).",
    "brainreg, brainregister, and ClearMap2 are all designed around full "
    "brains. Forcing a partial sample through them distorts the tissue to "
    "fill the missing atlas regions, regardless of orientation correctness.",
    "Both brainreg and brainregister fundamentally wrap a third-party engine "
    "(NiftyReg / elastix). If you want the same engine without the full-brain "
    "assumption, drive the engine directly with your own parameter file.",
], size=12)

# ==========================================================================
# 6. LANDMARK / INTERACTIVE METHODS
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Landmark-driven & interactive methods",
          "User picks corresponding points, math computes a smooth warp that "
          "satisfies them.")

cards = [
    dict(name="Thin-Plate Splines (TPS)",
         cite="Bookstein 1989 (IEEE PAMI)",
         body="Closed-form spline that lands every landmark exactly. With "
              "few/clustered landmarks it extrapolates wildly; fit "
              "on RESIDUALS of a rigid (vol2atlas's `--tps`) keeps the rigid valid far "
              "from landmarks.",
         partial=True,
         repo="scipy.interpolate.RBFInterpolator"),
    dict(name="3D Slicer + FRW",
         cite="SlicerIGT extension",
         body="Polished GUI for landmark-based TPS / affine / RBF. Outlier flagging, "
              "snap-to-feature, manages hundreds of fiducials cleanly. Heavier "
              "than the vol2atlas landmarks napari panel but more capable.",
         partial=True,
         repo="slicer.org · slicerigt.org"),
    dict(name="BigWarp",
         cite="Bogovic et al. 2016 (bioRxiv)",
         body="Fiji plugin for interactive TPS on very large volumes (OME-Zarr, "
              "N5, HDF5 directly). The standard tool when you want dense, hand-"
              "placed landmarks on a multi-TB stack.",
         partial=True,
         repo="imagej.net/plugins/bigwarp"),
    dict(name="ABBA (Aligning Big Brains and Atlases)",
         cite="BIOP, EPFL",
         body="Interactive whole-slide → CCF for 2D histology slices. Not a "
              "volume registration tool, but the spirit (manual prealign + "
              "elastix refine) and UI design ideas are relevant.",
         partial=False,
         repo="biop.epfl.ch/Image_ABBA.html"),
]
x0, y0 = Inches(0.4), Inches(1.5)
for i, c in enumerate(cards):
    col = i % 4
    method_card(s, x0 + col * (w + gap), y0, w, h, **c)

# ==========================================================================
# 7. SLICE-BASED TOOLS (NOT VOLUMES)
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Slice-based tools (2D histology → CCF)",
          "These appear in the same literature but solve a different problem — "
          "single coronal/sagittal slices, not 3D volumes.")
bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5), [
    ("DeepSlice.",
     "Carey et al. 2023 — CNN predicts the CCF coronal cut + warp for a single "
     "histology slice. Outputs an angle + position in CCFv3."),
    ("QUINT workflow.",
     "Yates et al. 2019 — three Fiji plugins (QuickNII, VisuAlign, Nutil) that "
     "manually align 2D histology to CCF for cell counting."),
    ("WholeBrain.",
     "Fürth et al. 2018 — R package, 2D histology → CCF with interactive "
     "outline matching and per-region counting."),
    ("ABBA.",
     "Mentioned on previous slide. 2D version of the same idea."),
    ("AMaSiNe.",
     "Sound et al. 2021 — adapts ANTs registration to 2D slices with intermediate "
     "3D consistency constraints."),
    ("Relevant for you? Mostly NO.",
     "Your input is a 3D volume; these are 2D-slice tools. They are listed here "
     "so you can recognise them when they come up in the mouse-brain registration "
     "literature and not waste time on a wrong fit."),
], size=14)

# ==========================================================================
# 8. COMPARISON TABLE
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "At-a-glance comparison",
          "Focused on what matters for partial-hemisphere µCT.")

# Build a manual table using rectangles + text
HEAD = ["Method", "Engine", "Partial sample", "Cross-modality", "CLI / API"]
ROWS = [
    ("vol2atlas rigid (prealign→refine→landmarks)",  "scipy",       "✓",  "n/a",  "CLI"),
    ("vol2atlas export --tps",     "RBFInterp",   "✓",  "n/a",  "CLI"),
    ("vol2atlas ants (SyNOnly)",   "antspyx",     "✓*", "fair", "CLI"),
    ("vol2atlas brainreg",         "NiftyReg",    "✗",  "fair", "CLI"),
    ("mBrainAligner (manual)","C++/Qt+CUDA", "✓",  "good", "binaries (legacy deps)"),
    ("3D Slicer + FRW",       "ITK + TPS",   "✓",  "n/a",  "GUI"),
    ("BigWarp (Fiji)",        "ImgLib2 TPS", "✓",  "n/a",  "GUI"),
    ("elastix custom",        "elastix",     "✓",  "fair", "CLI / Py"),
    ("VoxelMorph / SynthMorph","PyTorch",    "✗",  "good", "Py (needs training)"),
]

x0, y0 = Inches(0.45), Inches(1.55)
cell_w = [Inches(2.6), Inches(2.0), Inches(2.0), Inches(2.2), Inches(3.4)]
row_h  = Inches(0.42)

# header row
x = x0
for i, txt in enumerate(HEAD):
    box = add_panel(s, x, y0, cell_w[i], row_h, fill=ACCENT)
    add_text(s, x, y0 + Inches(0.07), cell_w[i], row_h,
             txt, size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    x += cell_w[i]

for r, row in enumerate(ROWS):
    y = y0 + (r + 1) * row_h
    bg_fill = PANEL if r % 2 == 0 else RGBColor(0x18, 0x18, 0x1f)
    x = x0
    for i, txt in enumerate(row):
        add_panel(s, x, y, cell_w[i], row_h, fill=bg_fill)
        # Color the "Partial sample" column
        col = FG
        if i == 2:
            col = OK if "✓" in txt else WARN
        add_text(s, x + Inches(0.05), y + Inches(0.08),
                 cell_w[i] - Inches(0.1), row_h,
                 txt, size=11, color=col, align=PP_ALIGN.CENTER,
                 mono=(i in (1, 4)))
        x += cell_w[i]

add_text(s, Inches(0.45), y0 + (len(ROWS) + 1.2) * row_h + Inches(0.1),
         Inches(12.5), Inches(0.4),
         "* vol2atlas ants uses masks built from the data; works in practice but "
         "the algorithm itself has no partial-sample concept — heavy "
         "regularization is required.", size=10, color=MUTED)

# ==========================================================================
# 9. RECOMMENDED PATHS
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Recommended paths for your data")

# Three lanes
lanes = [
    ("Conservative",
     OK,
     ["Stick with `vol2atlas export` (rigid + crop) as the final.",
      "Add ≥10 well-spread landmarks in `landmarks`, then `--tps`.",
      "Acceptable when downstream task is region lookup or rough "
      "overlay (~0.4 mm error is fine for atlas IDs)."]),
    ("Pragmatic",
     ACCENT,
     ["`vol2atlas ants out/rigid` with SyNOnly + auto-masks.",
      "Intensity-driven so it can pick up structure between landmarks.",
      "5-30 min runtime. Always compare to rigid in QC viewer "
      "— accept only if it's clearly better."]),
    ("Ambitious",
     YELLOW,
     ["More landmarks (15-30 well-spread) + `vol2atlas export --tps`.",
      "Or histogram-match the sample to the atlas, then aggressive ANTs "
      "(`--iter-high 100 --flow-sigma 3 --total-sigma 0`).",
      "mBrainAligner (Qu et al.) targets your exact case but its binary "
      "release needs Qt4 + OpenCV 3.4 + CUDA 10.0 — install hell on "
      "modern systems."]),
]
w = (SLIDE_W - Inches(1.0)) / 3 - Inches(0.15)
for i, (title, color, body) in enumerate(lanes):
    x = Inches(0.4) + i * (w + Inches(0.15))
    add_panel(s, x, Inches(1.6), w, Inches(5.5), border=color)
    add_text(s, x + Inches(0.2), Inches(1.75), w - Inches(0.4), Inches(0.5),
             title, size=20, bold=True, color=color)
    bullets(s, x + Inches(0.2), Inches(2.35), w - Inches(0.4), Inches(4.5),
            body, size=13, lead=color)

# ==========================================================================
# 10. KEY REFERENCES
# ==========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Key references")

refs = [
    ("ANTs SyN.",
     "Avants BB et al. (2008). Symmetric diffeomorphic image registration with "
     "cross-correlation. Med Image Anal 12(1):26-41."),
    ("elastix.",
     "Klein S et al. (2010). elastix: a toolbox for intensity-based medical "
     "image registration. IEEE TMI 29(1):196-205."),
    ("NiftyReg.",
     "Modat M et al. (2010). Fast free-form deformation using GPUs. "
     "Comput Methods Programs Biomed 98(3):278-284."),
    ("VoxelMorph.",
     "Balakrishnan G et al. (2019). VoxelMorph: a learning framework for "
     "deformable medical image registration. IEEE TMI 38(8):1788-1800."),
    ("SynthMorph.",
     "Hoffmann M et al. (2022). SynthMorph: learning contrast-invariant "
     "registration without acquired images. IEEE TMI 41(3):543-558."),
    ("mBrainAligner.",
     "Qu L et al. (2022). Cross-modal coherent registration of whole mouse "
     "brains. Nat Methods 19:111-118."),
    ("BrainGlobe brainreg.",
     "Tyson AL et al. (2022). A deep learning algorithm for 3D cell detection "
     "in whole mouse brain image datasets. PLoS Comp Biol; brainreg in BrainGlobe "
     "tools paper, Sci Reports 12:867."),
    ("BigWarp.",
     "Bogovic JA et al. (2016). Robust registration of calcium images by "
     "learned contrast synthesis. bioRxiv."),
    ("Bookstein TPS.",
     "Bookstein FL (1989). Principal warps: Thin-plate splines and the "
     "decomposition of deformations. IEEE PAMI 11(6):567-585."),
]
bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.7),
        refs, size=12)

# ==========================================================================
prs.save(str(OUT))
print(f"wrote {OUT}  ({OUT.stat().st_size/1e3:.1f} KB, "
      f"{len(prs.slides)} slides)")
